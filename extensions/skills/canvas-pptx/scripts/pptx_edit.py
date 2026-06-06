"""Declarative patch applier for PPTX editing.

Applies a JSON patch (from generate_edit_patch) to a PPTX file.
The agent never imports python-pptx directly -- this module handles
all shape-type-aware operations internally.

Usage:
    from pptx_edit import apply_edit_patch
    apply_edit_patch('input.pptx', patch, 'output.pptx')
"""

import copy
import csv
import io
import os
import re

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Inches


# Valid operations per shape type -- prevents AttributeError classes
_VALID_OPS = {
    "TEXT": {"set_text", "set_font_name", "set_font_size", "set_font_bold", "set_font_italic", "set_font_color"},
    "TEXT_BULLETS": {"set_text", "set_font_name", "set_font_size", "set_font_bold", "set_font_italic", "set_font_color"},
    "CHART": {"set_categories", "set_series", "set_chart_from"},
    "TABLE": {"set_data", "set_table_from"},
    "IMAGE": {"set_image"},
    "GROUP": set(),  # No direct edits on groups
}


class PatchError(Exception):
    """Raised when a patch operation fails with an actionable message."""
    pass


# Namespace for XML operations
_A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


# ---------------------------------------------------------------------------
# Source file readers ("Reference, Don't Transcribe")
# ---------------------------------------------------------------------------

def _parse_markdown_tables(text):
    """Parse pipe-delimited markdown tables from text.

    Returns list of tables, each table is list of rows (list of strings).
    """
    tables = []
    current_table = []
    in_table = False

    for line in text.split('\n'):
        stripped = line.strip()
        if '|' in stripped and not stripped.startswith('```'):
            cells = [c.strip() for c in stripped.split('|')]
            # Remove empty first/last from leading/trailing pipes
            if cells and cells[0] == '':
                cells = cells[1:]
            if cells and cells[-1] == '':
                cells = cells[:-1]

            # Skip separator rows (---|---|---)
            if all(re.match(r'^[-:]+$', c) for c in cells):
                continue

            if not in_table:
                in_table = True
                current_table = []
            current_table.append(cells)
        else:
            if in_table and current_table:
                tables.append(current_table)
                current_table = []
                in_table = False

    if in_table and current_table:
        tables.append(current_table)

    return tables


def _parse_csv_tables(text):
    """Parse CSV text into a single table (list of rows)."""
    reader = csv.reader(io.StringIO(text))
    return [list(row) for row in reader]


def _read_source_table(source_path, table_index, workspace_dir=None):
    """Read a table from a source file.

    Args:
        source_path: Path to source file (md, csv, xlsx)
        table_index: Which table to extract (0-based)
        workspace_dir: Base directory for resolving relative paths

    Returns:
        List of rows (list of strings), including header row.
    """
    if workspace_dir and not os.path.isabs(source_path):
        source_path = os.path.join(workspace_dir, source_path)

    if not os.path.exists(source_path):
        raise PatchError(f"Source file not found: {source_path}")

    ext = os.path.splitext(source_path)[1].lower()

    if ext in ('.md', '.markdown', '.txt'):
        with open(source_path, 'r', encoding='utf-8') as f:
            text = f.read()
        tables = _parse_markdown_tables(text)
        if table_index >= len(tables):
            raise PatchError(
                f"Table index {table_index} out of range. "
                f"File has {len(tables)} table(s)."
            )
        return tables[table_index]

    elif ext == '.csv':
        with open(source_path, 'r', encoding='utf-8') as f:
            text = f.read()
        table = _parse_csv_tables(text)
        if not table:
            raise PatchError(f"No data found in CSV file: {source_path}")
        return table

    elif ext in ('.xlsx', '.xls'):
        try:
            import openpyxl
        except ImportError:
            raise PatchError("openpyxl required to read Excel files. Install with: pip install openpyxl")
        wb = openpyxl.load_workbook(source_path, data_only=True)
        ws = wb.active
        table = []
        for row in ws.iter_rows(values_only=True):
            table.append([str(c) if c is not None else '' for c in row])
        if not table:
            raise PatchError(f"No data found in Excel file: {source_path}")
        return table

    else:
        raise PatchError(f"Unsupported source file type: {ext}")


def _extract_chart_data_from_source(ref, workspace_dir=None):
    """Extract chart categories and series from a source file reference.

    Args:
        ref: dict with keys: source, table_index, category_column, series_columns
        workspace_dir: Base directory for resolving relative paths

    Returns:
        (categories, series_dict) where categories is list of strings and
        series_dict maps series_name -> list of floats.
    """
    table = _read_source_table(ref['source'], ref.get('table_index', 0), workspace_dir)
    if len(table) < 2:
        raise PatchError(f"Table has only {len(table)} row(s), need at least header + 1 data row.")

    header = table[0]
    data_rows = table[1:]

    cat_col = ref['category_column']
    if cat_col not in header:
        raise PatchError(
            f"Column '{cat_col}' not found. Available: {header}"
        )
    cat_idx = header.index(cat_col)
    categories = [row[cat_idx] for row in data_rows]

    series_dict = {}
    for series_name, col_name in ref['series_columns'].items():
        if col_name not in header:
            raise PatchError(
                f"Column '{col_name}' not found. Available: {header}"
            )
        col_idx = header.index(col_name)
        values = []
        for row_num, row in enumerate(data_rows, 2):
            raw = row[col_idx] if col_idx < len(row) else ''
            values.append(_parse_number(raw, col_name, row_num))
        series_dict[series_name] = values

    return categories, series_dict


def _extract_table_data_from_source(ref, workspace_dir=None):
    """Extract table data from a source file reference.

    Args:
        ref: dict with keys: source, table_index, columns (optional)
        workspace_dir: Base directory for resolving relative paths

    Returns:
        List of rows (list of strings), including header row.
    """
    table = _read_source_table(ref['source'], ref.get('table_index', 0), workspace_dir)

    columns = ref.get('columns')
    if not columns:
        return table

    header = table[0]
    col_indices = []
    for col in columns:
        if col not in header:
            raise PatchError(
                f"Column '{col}' not found. Available: {header}"
            )
        col_indices.append(header.index(col))

    filtered = []
    for row in table:
        filtered.append([row[i] if i < len(row) else '' for i in col_indices])
    return filtered


def _parse_number(raw, col_name, row_num):
    """Parse a string value to float, stripping common formatting."""
    if isinstance(raw, (int, float)):
        return float(raw)
    cleaned = str(raw).strip()
    cleaned = cleaned.replace('$', '').replace(',', '').replace('%', '').strip()
    if not cleaned or cleaned == '-' or cleaned.lower() == 'n/a':
        return 0.0
    # Handle M/K/B suffixes
    multiplier = 1.0
    if cleaned[-1].upper() == 'M':
        multiplier = 1_000_000
        cleaned = cleaned[:-1]
    elif cleaned[-1].upper() == 'K':
        multiplier = 1_000
        cleaned = cleaned[:-1]
    elif cleaned[-1].upper() == 'B':
        multiplier = 1_000_000_000
        cleaned = cleaned[:-1]
    try:
        return float(cleaned) * multiplier
    except ValueError:
        raise PatchError(
            f"Could not parse '{raw}' as number in column '{col_name}', row {row_num}"
        )


# ---------------------------------------------------------------------------
# Format-preserving text replacement
# ---------------------------------------------------------------------------

def _snapshot_run_format(run):
    """Capture formatting from a run for later restoration.

    Captures both the python-pptx API properties AND the raw rPr XML element
    to preserve properties not exposed by the API (e.g., latin/ea/cs typeface
    elements, solidFill with scheme colors, highlight, strikethrough, etc.).
    """
    fmt = {}
    font = run.font

    fmt['bold'] = font.bold
    fmt['italic'] = font.italic
    fmt['underline'] = font.underline
    fmt['size'] = font.size
    fmt['name'] = font.name

    # Safely capture color
    try:
        if font.color and font.color.rgb:
            fmt['color_rgb'] = font.color.rgb
        elif font.color and font.color.theme_color:
            fmt['color_theme'] = font.color.theme_color
    except AttributeError:
        pass

    # Capture full rPr XML for faithful restoration
    _a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    rPr = run._r.find(f'{{{_a_ns}}}rPr')
    if rPr is not None:
        fmt['_rPr_xml'] = copy.deepcopy(rPr)

    return fmt


def _apply_run_format(run, fmt):
    """Apply previously captured formatting to a run.

    If a full rPr XML snapshot is available, replaces the run's rPr element
    entirely to preserve all formatting faithfully (fonts, colors, effects).
    """
    _a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # If we have full rPr XML, use it for faithful restoration
    rPr_xml = fmt.get('_rPr_xml')
    if rPr_xml is not None:
        existing_rPr = run._r.find(f'{{{_a_ns}}}rPr')
        new_rPr = copy.deepcopy(rPr_xml)
        if existing_rPr is not None:
            run._r.replace(existing_rPr, new_rPr)
        else:
            run._r.insert(0, new_rPr)
        return

    # Fallback: apply individual properties via API
    if fmt.get('bold') is not None:
        run.font.bold = fmt['bold']
    if fmt.get('italic') is not None:
        run.font.italic = fmt['italic']
    if fmt.get('underline') is not None:
        run.font.underline = fmt['underline']
    if fmt.get('size') is not None:
        run.font.size = fmt['size']
    if fmt.get('name') is not None:
        run.font.name = fmt['name']
    if 'color_rgb' in fmt:
        run.font.color.rgb = fmt['color_rgb']
    elif 'color_theme' in fmt:
        run.font.color.theme_color = fmt['color_theme']


def _snapshot_paragraph_format(paragraph):
    """Capture paragraph-level formatting including full XML pPr element.

    The pPr element contains bullet definitions (buFont, buChar, buAutoNum),
    margins (marL), indents, line spacing, and other properties that the
    python-pptx API doesn't fully expose. Capturing the raw XML ensures
    faithful format preservation when rewriting bullet content.
    """
    fmt = {}

    fmt['alignment'] = paragraph.alignment
    fmt['level'] = paragraph.level

    if paragraph.space_before is not None:
        fmt['space_before'] = paragraph.space_before
    if paragraph.space_after is not None:
        fmt['space_after'] = paragraph.space_after
    if paragraph.line_spacing is not None:
        fmt['line_spacing'] = paragraph.line_spacing

    # Capture full pPr XML for faithful restoration (bullets, margins, indents)
    _a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    pPr = paragraph._p.find(f'{{{_a_ns}}}pPr')
    if pPr is not None:
        fmt['_pPr_xml'] = copy.deepcopy(pPr)

    # Snapshot run format from first run if available
    if paragraph.runs:
        fmt['run_format'] = _snapshot_run_format(paragraph.runs[0])

    return fmt


def _apply_paragraph_format(paragraph, fmt):
    """Apply previously captured paragraph formatting.

    If a full pPr XML snapshot is available, replaces the paragraph's pPr
    element entirely to preserve bullets, margins, indents, and other
    properties not exposed by the python-pptx API.
    """
    _a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # If we have a full pPr XML snapshot, use it for faithful restoration
    pPr_xml = fmt.get('_pPr_xml')
    if pPr_xml is not None:
        existing_pPr = paragraph._p.find(f'{{{_a_ns}}}pPr')
        new_pPr = copy.deepcopy(pPr_xml)
        if existing_pPr is not None:
            paragraph._p.replace(existing_pPr, new_pPr)
        else:
            # Insert pPr as first child of <a:p>
            paragraph._p.insert(0, new_pPr)
        return

    # Fallback: apply individual properties via API
    if fmt.get('alignment') is not None:
        paragraph.alignment = fmt['alignment']
    if fmt.get('level') is not None:
        paragraph.level = fmt['level']

    if 'space_before' in fmt:
        paragraph.space_before = fmt['space_before']
    if 'space_after' in fmt:
        paragraph.space_after = fmt['space_after']
    if 'line_spacing' in fmt:
        paragraph.line_spacing = fmt['line_spacing']


# ---------------------------------------------------------------------------
# Shape operation handlers
# ---------------------------------------------------------------------------

def _apply_text(shape, new_text):
    """Replace text in a shape, preserving formatting from the first run."""
    tf = shape.text_frame

    if isinstance(new_text, list):
        _apply_bullets(shape, new_text)
        return

    # Snapshot formatting from existing content
    run_fmt = None
    para_fmt = None
    if tf.paragraphs:
        para_fmt = _snapshot_paragraph_format(tf.paragraphs[0])
        if tf.paragraphs[0].runs:
            run_fmt = _snapshot_run_format(tf.paragraphs[0].runs[0])

    # Clear all paragraphs
    tf.clear()

    # Set new text
    p = tf.paragraphs[0]
    if para_fmt:
        _apply_paragraph_format(p, para_fmt)

    run = p.add_run()
    run.text = str(new_text)
    if run_fmt:
        _apply_run_format(run, run_fmt)
    elif para_fmt and 'run_format' in para_fmt:
        _apply_run_format(run, para_fmt['run_format'])


def _apply_bullets(shape, items):
    """Replace bullet text in a shape, preserving formatting per-paragraph."""
    tf = shape.text_frame

    # Snapshot formatting from existing paragraphs
    para_fmts = []
    for p in tf.paragraphs:
        para_fmts.append(_snapshot_paragraph_format(p))

    # Use first paragraph's format as default for all
    default_fmt = para_fmts[0] if para_fmts else {}

    # Clear
    tf.clear()

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        fmt = para_fmts[i] if i < len(para_fmts) else default_fmt
        _apply_paragraph_format(p, fmt)

        run = p.add_run()
        run.text = str(item)
        if 'run_format' in fmt:
            _apply_run_format(run, fmt['run_format'])


def _apply_chart_data(shape, categories, series_dict):
    """Replace chart data using CategoryChartData."""
    chart = shape.chart
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_dict.items():
        chart_data.add_series(name, values)
    chart.replace_data(chart_data)


def _apply_table_data(shape, data):
    """Replace table cell contents, preserving cell formatting."""
    table = shape.table
    rows_needed = len(data)
    cols_needed = len(data[0]) if data else 0

    for r in range(min(rows_needed, len(table.rows))):
        for c in range(min(cols_needed, len(table.columns))):
            cell = table.cell(r, c)
            # Preserve formatting from existing cell
            run_fmt = None
            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                run_fmt = _snapshot_run_format(cell.text_frame.paragraphs[0].runs[0])

            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(data[r][c])
            if run_fmt:
                _apply_run_format(run, run_fmt)


def _apply_image(shape, image_path, slide, workspace_dir=None):
    """Replace an image while preserving position and size."""
    if workspace_dir and not os.path.isabs(image_path):
        image_path = os.path.join(workspace_dir, image_path)

    if not os.path.exists(image_path):
        raise PatchError(f"Image file not found: {image_path}")

    # Preserve position/size
    left, top, width, height = shape.left, shape.top, shape.width, shape.height

    # Remove old shape's element from slide
    sp = shape._element
    sp.getparent().remove(sp)

    # Add new image at same position
    slide.shapes.add_picture(image_path, left, top, width, height)


def _apply_font_style(shape, shape_spec):
    """Apply font style changes (name, size, bold, italic, color) to all runs in a shape."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    if not shape.has_text_frame:
        return

    font_name = shape_spec.get('set_font_name')
    font_size = shape_spec.get('set_font_size')
    font_bold = shape_spec.get('set_font_bold')
    font_italic = shape_spec.get('set_font_italic')
    font_color = shape_spec.get('set_font_color')

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if font_name is not None:
                run.font.name = font_name
            if font_size is not None:
                run.font.size = Pt(font_size)
            if font_bold is not None:
                run.font.bold = font_bold
            if font_italic is not None:
                run.font.italic = font_italic
            if font_color is not None:
                hex_val = font_color.lstrip('#')
                r = int(hex_val[0:2], 16)
                g = int(hex_val[2:4], 16)
                b = int(hex_val[4:6], 16)
                run.font.color.rgb = RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Slide-level color replacement
# ---------------------------------------------------------------------------

def _apply_color_replace(slide, color_map):
    """Replace colors across all shapes on a slide.

    Args:
        slide: python-pptx slide object
        color_map: dict mapping old_hex -> new_hex (e.g., {"4338CA": "EA580C"})
                   Hex values are 6-char uppercase, no '#' prefix.
    """
    from lxml import etree

    # Normalize the map to uppercase
    normalized = {k.upper().lstrip('#'): v.upper().lstrip('#') for k, v in color_map.items()}

    # Walk all elements and replace color val attributes in-place
    for elem in slide._element.iter():
        val = elem.get('val')
        if val and len(val) == 6:
            upper_val = val.upper()
            if upper_val in normalized:
                elem.set('val', normalized[upper_val])
        # Also check lastClr attribute (system colors)
        last_clr = elem.get('lastClr')
        if last_clr and len(last_clr) == 6:
            upper_lc = last_clr.upper()
            if upper_lc in normalized:
                elem.set('lastClr', normalized[upper_lc])


def _get_slide_colors(slide):
    """Extract all hex colors used in a slide (for generating color_replace suggestions).

    Returns dict of hex_color -> count.
    """
    from lxml import etree
    from collections import Counter

    colors = Counter()
    slide_xml = etree.tostring(slide._element, encoding='unicode')

    # Find all srgbClr val="XXXXXX"
    import re
    for match in re.finditer(r'val="([0-9A-Fa-f]{6})"', slide_xml):
        colors[match.group(1).upper()] += 1

    return dict(colors.most_common(20))


# ---------------------------------------------------------------------------
# Slide addition (hybrid approach)
# ---------------------------------------------------------------------------

def _get_style_from_slide(prs, slide_num, hint='body'):
    """Extract font style from a reference slide's shape.

    Args:
        prs: Presentation object
        slide_num: 1-based slide number to extract style from
        hint: 'title' to find the title shape, 'body' to find the largest text shape

    Returns dict with font properties, or None if not found.
    """
    if slide_num < 1 or slide_num > len(prs.slides):
        return None
    slide = prs.slides[slide_num - 1]
    best = None
    best_area = 0
    for shape in slide.shapes:
        if not shape.has_text_frame or not shape.text_frame.paragraphs:
            continue
        if hint == 'title' and shape == slide.shapes.title:
            best = shape
            break
        area = shape.width * shape.height
        if hint == 'body' and shape != slide.shapes.title and area > best_area:
            best = shape
            best_area = area
    if not best or not best.text_frame.paragraphs[0].runs:
        return None
    run = best.text_frame.paragraphs[0].runs[0]
    style = {}
    if run.font.name:
        style['name'] = run.font.name
    if run.font.size:
        style['size'] = run.font.size
    if run.font.bold is not None:
        style['bold'] = run.font.bold
    if run.font.italic is not None:
        style['italic'] = run.font.italic
    if run.font.color and run.font.color.type is not None:
        try:
            style['color_rgb'] = run.font.color.rgb
        except Exception:
            pass
    return style if style else None


def _create_textbox_fallback(slide, name, value, position_hint=None, ref_style=None,
                             slide_width=None, slide_height=None):
    """Create a textbox as fallback when a placeholder is missing.

    Uses position_hint ('title' or 'body') to pick sensible defaults.
    If ref_style is provided (from _get_style_from_slide), uses those
    font properties instead of hardcoded defaults.
    If slide_width/slide_height are provided, scales positions to match
    the actual slide dimensions instead of assuming 10" x 7.5".
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # Base dimensions assume standard 10" x 7.5" — scale if actual dims differ
    sw = slide_width if slide_width else Inches(10)
    sh = slide_height if slide_height else Inches(7.5)
    # Scale factors relative to 10" x 7.5" base
    sx = sw / Inches(10)
    sy = sh / Inches(7.5)

    if position_hint == 'title':
        left = int(Inches(0.5) * sx)
        top = int(Inches(0.3) * sy)
        width = int(Inches(9) * sx)
        height = int(Inches(1) * sy)
        font_size, font_bold = Pt(28), True
    else:
        left = int(Inches(0.5) * sx)
        top = int(Inches(1.5) * sy)
        width = int(Inches(9) * sx)
        height = int(Inches(5.5) * sy)
        font_size, font_bold = Pt(14), False

    # Override defaults with reference style if available
    if ref_style:
        font_size = ref_style.get('size', font_size)
        font_bold = ref_style.get('bold', font_bold)

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    def _apply_ref_style(run):
        run.font.size = font_size
        run.font.bold = font_bold
        if ref_style:
            if 'name' in ref_style:
                run.font.name = ref_style['name']
            if 'italic' in ref_style:
                run.font.italic = ref_style['italic']
            if 'color_rgb' in ref_style:
                run.font.color.rgb = ref_style['color_rgb']

    if isinstance(value, list):
        for i, item in enumerate(value):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = str(item)
            _apply_ref_style(run)
            p.space_after = Pt(8)
    else:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(value)
        _apply_ref_style(run)

    return txbox


def _add_layout_based_slides(prs, add_slides_spec, workspace_dir=None):
    """Add new slides from layout + content specs.

    Returns list of (insert_at, slide) tuples.

    If the requested layout is not found, falls back to the first available
    layout. If a placeholder is missing (common in minimal-layout PPTX files
    with only 1-2 layouts), creates a textbox at a sensible default position
    instead of failing. If style_from_slide is set, clones font styling from
    that slide's shapes for the fallback textboxes.
    """
    new_slides = []
    layout_map = {layout.name: layout for layout in prs.slide_layouts}

    for spec in add_slides_spec:
        layout_name = spec.get('layout')
        if layout_name not in layout_map:
            available = list(layout_map.keys())
            # Fallback: use first available layout instead of failing
            if available:
                fallback = available[0]
                print(f"WARNING: Layout '{layout_name}' not found. "
                      f"Falling back to '{fallback}'. Available: {available}")
                layout_name = fallback
            else:
                raise PatchError("No slide layouts available in this presentation.")

        layout = layout_map[layout_name]
        slide = prs.slides.add_slide(layout)

        # Pre-fetch reference styles if style_from_slide is specified
        style_ref = spec.get('style_from_slide')
        title_style = None
        body_style = None
        if style_ref:
            title_style = _get_style_from_slide(prs, style_ref, hint='title')
            body_style = _get_style_from_slide(prs, style_ref, hint='body')
            if title_style or body_style:
                print(f"  Using style from slide {style_ref}: "
                      f"title={title_style}, body={body_style}")

        content = spec.get('content', {})
        content_items = list(content.items())
        for idx, (ph_name, value) in enumerate(content_items):
            # Find placeholder by name
            ph = None
            for p in slide.placeholders:
                if p.name == ph_name:
                    ph = p
                    break
            if ph is not None:
                _apply_text(ph, value)
            else:
                # Fallback: create a textbox instead of failing
                hint = 'title' if idx == 0 else 'body'
                ref_style = title_style if hint == 'title' else body_style
                print(f"WARNING: Placeholder '{ph_name}' not found in layout "
                      f"'{layout_name}'. Creating textbox fallback ({hint}).")
                _create_textbox_fallback(slide, ph_name, value,
                                         position_hint=hint, ref_style=ref_style,
                                         slide_width=prs.slide_width,
                                         slide_height=prs.slide_height)

        # Handle chart/table from source on new slides
        for shape_spec in spec.get('shapes', []):
            _apply_shape_ops_on_slide(slide, shape_spec, workspace_dir)

        insert_at = spec.get('insert_at', len(prs.slides))
        new_slides.append((insert_at, slide))

    return new_slides


def _apply_shape_ops_on_slide(slide, shape_spec, workspace_dir=None):
    """Apply operations to a named shape on a slide."""
    name = shape_spec.get('name')
    target = None
    for shape in slide.shapes:
        if shape.name == name:
            target = shape
            break
    if target is None:
        return

    shape_type = shape_spec.get('type', 'OTHER')
    _apply_shape_ops(target, shape_spec, shape_type, slide, workspace_dir)


# ---------------------------------------------------------------------------
# Main applier
# ---------------------------------------------------------------------------

def _validate_ops(shape_spec, shape_type):
    """Validate that requested operations are valid for the shape type."""
    valid = _VALID_OPS.get(shape_type, set())
    for key in shape_spec:
        if key.startswith('set_') and shape_spec[key] is not None:
            if key not in valid:
                raise PatchError(
                    f"Operation '{key}' is not valid for shape type '{shape_type}'. "
                    f"Valid operations: {valid or 'none'}"
                )


def _apply_shape_ops(shape, shape_spec, shape_type, slide, workspace_dir=None):
    """Apply all set_* operations from a shape spec to a shape."""
    _validate_ops(shape_spec, shape_type)

    # Text
    if shape_spec.get('set_text') is not None:
        _apply_text(shape, shape_spec['set_text'])

    # Chart from source reference
    if shape_spec.get('set_chart_from') is not None:
        ref = shape_spec['set_chart_from']
        categories, series_dict = _extract_chart_data_from_source(ref, workspace_dir)
        _apply_chart_data(shape, categories, series_dict)

    # Chart inline
    elif shape_spec.get('set_series') is not None:
        categories = shape_spec.get('set_categories', [])
        _apply_chart_data(shape, categories, shape_spec['set_series'])

    # Table from source reference
    if shape_spec.get('set_table_from') is not None:
        ref = shape_spec['set_table_from']
        data = _extract_table_data_from_source(ref, workspace_dir)
        _apply_table_data(shape, data)

    # Table inline
    elif shape_spec.get('set_data') is not None:
        _apply_table_data(shape, shape_spec['set_data'])

    # Image
    if shape_spec.get('set_image') is not None:
        _apply_image(shape, shape_spec['set_image'], slide, workspace_dir)

    # Font style changes
    font_ops = ('set_font_name', 'set_font_size', 'set_font_bold', 'set_font_italic', 'set_font_color')
    if any(shape_spec.get(op) is not None for op in font_ops):
        _apply_font_style(shape, shape_spec)


_PML_NS = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
_REL_NS = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'


def _get_sld_id_lst(prs):
    """Get the sldIdLst element from the presentation XML."""
    return prs.part._element.find(f'{_PML_NS}sldIdLst')


def _reorder_slides(prs, new_order):
    """Reorder slides according to new_order (1-based slide numbers)."""
    # Validate
    current_count = len(prs.slides)
    if sorted(new_order) != list(range(1, current_count + 1)):
        raise PatchError(
            f"reorder must contain every slide number exactly once. "
            f"Expected {list(range(1, current_count + 1))}, got {sorted(new_order)}"
        )

    # Access the XML sldIdLst
    sld_id_lst = _get_sld_id_lst(prs)
    sld_ids = list(sld_id_lst)

    # Clear and re-add in new order
    for sid in sld_ids:
        sld_id_lst.remove(sid)

    for num in new_order:
        sld_id_lst.append(sld_ids[num - 1])


def _delete_slides(prs, slide_numbers):
    """Delete slides by number (1-based). Deletes from highest to lowest."""
    for num in sorted(slide_numbers, reverse=True):
        idx = num - 1
        if idx < 0 or idx >= len(prs.slides):
            raise PatchError(f"Slide {num} does not exist (presentation has {len(prs.slides)} slides)")

        slide = prs.slides[idx]
        rId = None
        for rel in prs.part.rels.values():
            if rel.target_part == slide.part:
                rId = rel.rId
                break

        if rId:
            # Remove from sldIdLst
            sld_id_lst = _get_sld_id_lst(prs)
            for sld_id in sld_id_lst:
                if sld_id.get(f'{_REL_NS}id') == rId:
                    sld_id_lst.remove(sld_id)
                    break
            # Remove relationship
            prs.part.rels.pop(rId)


def _clone_slide_from_template(prs, clone_from_idx):
    """Create a new slide by deep-copying a source slide within the same presentation.

    Unlike add_slide(layout), this preserves all visual elements from the source:
    backgrounds, decorative shapes, branded bars, images, etc.

    Args:
        prs: The open Presentation object.
        clone_from_idx: 0-based index of the slide to clone.

    Returns:
        The newly appended slide (with full visual content from source).
    """
    _PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    # Add a structural stub using the last layout (cheapest skeleton)
    base_layout = prs.slide_layouts[len(prs.slide_layouts) - 1]
    new_slide = prs.slides.add_slide(base_layout)

    src_slide = prs.slides[clone_from_idx]

    # Replace the new slide's cSld with a deep copy of the source's cSld.
    # cSld contains spTree (all shapes), background fill, and notes references —
    # everything needed to reproduce the slide's visual appearance.
    src_csld = src_slide._element.find(f'{{{_PML}}}cSld')
    new_csld = new_slide._element.find(f'{{{_PML}}}cSld')
    if src_csld is not None and new_csld is not None:
        new_slide._element.replace(new_csld, copy.deepcopy(src_csld))

    # Mirror image relationships so embedded pictures resolve correctly.
    # We must preserve the original rId because the deep-copied cSld XML contains
    # r:embed="rIdN" attributes that reference specific rIds by name.
    for rel in src_slide.part.rels.values():
        if 'image' in rel.reltype and not getattr(rel, 'is_external', False):
            try:
                # Insert directly into _rels with the same rId so XML references resolve
                new_slide.part.rels._rels[rel.rId] = rel
            except Exception:
                pass

    return new_slide


def _move_slide_to_position(prs, slide, target_pos):
    """Move a slide (already appended at end) to target_pos (1-based)."""
    sld_id_lst = _get_sld_id_lst(prs)
    sld_ids = list(sld_id_lst)

    if len(sld_ids) < 2:
        return

    # The new slide is at the end
    last = sld_ids[-1]

    # Target index (0-based, clamped)
    target_idx = max(0, min(target_pos - 1, len(sld_ids) - 1))

    if target_idx == len(sld_ids) - 1:
        return  # Already at correct position

    # Remove from end
    sld_id_lst.remove(last)

    # Insert at target position
    sld_ids_updated = list(sld_id_lst)
    sld_ids_updated.insert(target_idx, last)

    # Rebuild
    for sid in list(sld_id_lst):
        sld_id_lst.remove(sid)
    for sid in sld_ids_updated:
        sld_id_lst.append(sid)


def apply_edit_patch(input_path, patch, output_path, workspace_dir=None):
    """Apply a declarative edit patch to a PPTX file.

    Args:
        input_path: Path to the source PPTX file
        patch: Dict from generate_edit_patch() with agent's modifications
        output_path: Path to write the modified PPTX
        workspace_dir: Base directory for resolving relative paths in source references

    Returns:
        Dict with summary of operations performed.
    """
    if workspace_dir is None:
        workspace_dir = os.path.dirname(os.path.abspath(input_path))

    prs = Presentation(input_path)
    summary = {"modified_shapes": 0, "added_slides": 0, "deleted_slides": 0, "errors": []}

    # Capture sld_id elements for slides to delete BEFORE any additions shift positions.
    # delete_slides numbers are 1-based positions in the original input file.
    structure = patch.get('structure', {})
    delete_list = structure.get('delete_slides', [])
    sld_id_elements_to_delete = []
    if delete_list:
        snapshot = list(_get_sld_id_lst(prs))
        for num in delete_list:
            idx = num - 1
            if 0 <= idx < len(snapshot):
                sld_id_elements_to_delete.append(snapshot[idx])
            else:
                summary["errors"].append(
                    f"delete_slides: slide {num} out of range "
                    f"(template has {len(snapshot)} slides)"
                )

    # Step 1: Apply modifications to existing slides
    for slide_spec in patch.get('slides', []):
        slide_num = slide_spec.get('slide')
        if slide_num is None:
            continue

        idx = slide_num - 1
        if idx < 0 or idx >= len(prs.slides):
            summary["errors"].append(f"Slide {slide_num} does not exist")
            continue

        slide = prs.slides[idx]

        for shape_spec in slide_spec.get('shapes', []):
            shape_name = shape_spec.get('name')
            shape_type = shape_spec.get('type', 'OTHER')

            # Find shape by name
            target = None
            for shape in slide.shapes:
                if shape.name == shape_name:
                    target = shape
                    break

            if target is None:
                summary["errors"].append(f"Slide {slide_num}: shape '{shape_name}' not found")
                continue

            # Check if any set_* field is non-null
            has_changes = any(
                k.startswith('set_') and v is not None
                for k, v in shape_spec.items()
            )
            if not has_changes:
                continue

            try:
                _apply_shape_ops(target, shape_spec, shape_type, slide, workspace_dir)
                summary["modified_shapes"] += 1
            except PatchError as e:
                summary["errors"].append(f"Slide {slide_num}, '{shape_name}': {e}")
            except Exception as e:
                summary["errors"].append(f"Slide {slide_num}, '{shape_name}': unexpected error: {e}")

        # Slide-level color replacement
        color_map = slide_spec.get('color_replace')
        if color_map:
            try:
                _apply_color_replace(slide, color_map)
                summary["modified_shapes"] += 1  # count as a modification
            except Exception as e:
                summary["errors"].append(f"Slide {slide_num} color_replace: {e}")

    # Steps 2 & 2b: Add new slides (layout-based and template clones).
    # All new slides are first appended at the end, then repositioned together
    # in a single pass so that mixing clone_slides and add_slides in one patch
    # produces the correct final order.
    original_count = len(prs.slides)
    new_slide_assignments = []  # (sld_id_element, insert_at)

    # Step 2: Layout-based new slides (add_slides)
    for spec in patch.get('add_slides', []):
        try:
            layout_name = spec.get('layout')
            layout_map = {layout.name: layout for layout in prs.slide_layouts}
            if layout_name not in layout_map:
                available = list(layout_map.keys())
                if available:
                    fallback = available[0]
                    print(f"WARNING: Layout '{layout_name}' not found. "
                          f"Falling back to '{fallback}'. Available: {available}")
                    layout_name = fallback
                else:
                    raise PatchError("No slide layouts available in this presentation.")

            layout = layout_map[layout_name]
            slide = prs.slides.add_slide(layout)
            sld_id = list(_get_sld_id_lst(prs))[-1]  # sld_id element for this slide

            # Pre-fetch reference styles if style_from_slide is specified
            style_ref = spec.get('style_from_slide')
            title_style = None
            body_style = None
            if style_ref:
                title_style = _get_style_from_slide(prs, style_ref, hint='title')
                body_style = _get_style_from_slide(prs, style_ref, hint='body')
                if title_style or body_style:
                    print(f"  Using style from slide {style_ref}: "
                          f"title={title_style}, body={body_style}")

            content = spec.get('content', {})
            content_items = list(content.items())
            for idx, (ph_name, value) in enumerate(content_items):
                ph = None
                for p in slide.placeholders:
                    if p.name == ph_name:
                        ph = p
                        break
                if ph is not None:
                    _apply_text(ph, value)
                else:
                    hint = 'title' if idx == 0 else 'body'
                    ref_style = title_style if hint == 'title' else body_style
                    print(f"WARNING: Placeholder '{ph_name}' not found in layout "
                          f"'{layout_name}'. Creating textbox fallback ({hint}).")
                    _create_textbox_fallback(slide, ph_name, value,
                                             position_hint=hint, ref_style=ref_style,
                                             slide_width=prs.slide_width,
                                             slide_height=prs.slide_height)

            insert_at = spec.get('insert_at', original_count + len(new_slide_assignments) + 1)
            new_slide_assignments.append((sld_id, insert_at))
            summary["added_slides"] += 1

        except PatchError as e:
            summary["errors"].append(f"add_slides: {e}")
        except Exception as e:
            summary["errors"].append(f"add_slides: unexpected error: {e}")

    # Step 2b: Clone template slides (full visual copy, preserves all design elements)
    for spec in patch.get('clone_slides', []):
        clone_from = spec.get('clone_from')
        if clone_from is None:
            summary["errors"].append("clone_slides: missing 'clone_from' field")
            continue
        src_idx = clone_from - 1  # convert to 0-based; references original slides
        if src_idx < 0 or src_idx >= original_count:
            summary["errors"].append(
                f"clone_slides: clone_from={clone_from} out of range "
                f"(template has {original_count} slides)"
            )
            continue
        try:
            new_slide = _clone_slide_from_template(prs, src_idx)
            sld_id = list(_get_sld_id_lst(prs))[-1]  # sld_id element for this slide

            # Apply text content replacements to named shapes/placeholders
            content = spec.get('content', {})
            for shape_name, value in content.items():
                matched = None
                for shape in new_slide.shapes:
                    if shape.name == shape_name:
                        matched = shape
                        break
                if matched is not None and matched.has_text_frame:
                    _apply_text(matched, value)
                else:
                    print(f"WARNING: clone_slides: shape '{shape_name}' not found "
                          f"(cloned from slide {clone_from})")

            insert_at = spec.get('insert_at', original_count + len(new_slide_assignments) + 1)
            new_slide_assignments.append((sld_id, insert_at))
            summary["added_slides"] += 1

        except Exception as e:
            summary["errors"].append(f"clone_slides: {e}")

    # Reposition all new slides in a single coordinated pass.
    # new_slide_assignments contains (sld_id, insert_at) in append order.
    # The original slides occupy positions 1..original_count; new slides are
    # inserted among them according to each insert_at value.
    if new_slide_assignments:
        sld_id_lst = _get_sld_id_lst(prs)
        all_ids = list(sld_id_lst)
        orig_ids = all_ids[:original_count]
        new_id_map = {id(sld_id): insert_at for sld_id, insert_at in new_slide_assignments}

        # Build final order: start from originals, insert each new slide at its target
        result = list(orig_ids)
        for sld_id, ins_pos in sorted(new_slide_assignments, key=lambda x: x[1]):
            idx = max(0, min(ins_pos - 1, len(result)))
            result.insert(idx, sld_id)

        # Rebuild sld_id_lst with the computed order
        for sid in list(sld_id_lst):
            sld_id_lst.remove(sid)
        for sid in result:
            sld_id_lst.append(sid)

    # Step 3: Delete original slides by element identity.
    # sld_id elements were captured before any new slides were added, so these
    # always reference the correct original slides regardless of repositioning.
    if sld_id_elements_to_delete:
        sld_id_lst = _get_sld_id_lst(prs)
        current_ids = list(sld_id_lst)
        for sld_id_elem in sld_id_elements_to_delete:
            if sld_id_elem in current_ids:
                rId = sld_id_elem.get(f'{_REL_NS}id')
                sld_id_lst.remove(sld_id_elem)
                if rId and rId in prs.part.rels:
                    prs.part.rels.pop(rId)
        summary["deleted_slides"] = len(sld_id_elements_to_delete)

    reorder = structure.get('reorder')
    if reorder:
        try:
            _reorder_slides(prs, reorder)
        except PatchError as e:
            summary["errors"].append(f"reorder: {e}")

    # Save
    prs.save(output_path)

    # Print summary
    print(f"Patch applied: {summary['modified_shapes']} shapes modified, "
          f"{summary['added_slides']} slides added, "
          f"{summary['deleted_slides']} slides deleted")
    if summary["errors"]:
        print(f"Errors ({len(summary['errors'])}):")
        for err in summary["errors"]:
            print(f"  - {err}")

    return summary


# ---------------------------------------------------------------------------
# Slide merging (insert slides from a separately-created PPTX)
# ---------------------------------------------------------------------------

def merge_slides(target_path, source_path, output_path, insert_at=None):
    """Merge slides from a source PPTX into a target PPTX.

    Use this after apply_edit_patch() to insert slides created via pptxgenjs
    (creation.md) into an edited presentation.

    Args:
        target_path: Path to the target PPTX (the edited presentation)
        source_path: Path to the source PPTX (created by pptxgenjs with new slides)
        output_path: Path to write the merged PPTX
        insert_at: 1-based position to insert the new slides.
                   None or 0 = append at end.

    Returns:
        Dict with summary of merge operation.
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from lxml import etree
    import copy as _copy

    target_prs = Presentation(target_path)
    source_prs = Presentation(source_path)

    summary = {"slides_merged": 0, "errors": []}

    for src_slide in source_prs.slides:
        try:
            # Use a blank layout from the target as the base
            # (the actual content comes from the source slide's XML)
            blank_layout = target_prs.slide_layouts[len(target_prs.slide_layouts) - 1]
            new_slide = target_prs.slides.add_slide(blank_layout)

            # Replace the new slide's XML with the source slide's XML
            # but keep the new slide's relationship infrastructure
            src_xml = etree.tostring(src_slide._element)
            new_elem = etree.fromstring(src_xml)

            # Preserve the new slide's spTree (shape tree) by replacing it
            src_sp_tree = src_slide._element.find(
                './/{http://schemas.openxmlformats.org/presentationml/2006/main}cSld/'
                '{http://schemas.openxmlformats.org/presentationml/2006/main}spTree'
            )
            if src_sp_tree is None:
                src_sp_tree = src_slide._element.find(
                    './/{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}spTree'
                )
            # Fallback: look for spTree with common namespace pattern
            if src_sp_tree is None:
                for elem in src_slide._element.iter():
                    if elem.tag.endswith('}spTree') or elem.tag == 'spTree':
                        src_sp_tree = elem
                        break

            # Copy the entire cSld element (contains all shapes, background, etc.)
            _PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            src_csld = src_slide._element.find(f'{{{_PML}}}cSld')
            new_csld = new_slide._element.find(f'{{{_PML}}}cSld')

            if src_csld is not None and new_csld is not None:
                # Replace the target cSld with source cSld
                parent = new_slide._element
                parent.replace(new_csld, _copy.deepcopy(src_csld))

            # Copy images and other media from source slide.
            # Preserve the original rId so r:embed attributes in the copied XML resolve.
            for rel in src_slide.part.rels.values():
                if "image" in rel.reltype and not getattr(rel, 'is_external', False):
                    try:
                        new_slide.part.rels._rels[rel.rId] = rel
                    except Exception:
                        pass

            # Move to insert position if specified
            if insert_at and insert_at > 0:
                _move_slide_to_position(target_prs, new_slide, insert_at)

            summary["slides_merged"] += 1

        except Exception as e:
            summary["errors"].append(f"merge slide: {e}")

    target_prs.save(output_path)

    print(f"Merge complete: {summary['slides_merged']} slides merged")
    if insert_at:
        print(f"  Inserted at position {insert_at}")
    if summary["errors"]:
        print(f"Errors ({len(summary['errors'])}):")
        for err in summary["errors"]:
            print(f"  - {err}")

    return summary
