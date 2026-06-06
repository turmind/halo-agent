"""Comprehensive template inspection utility for PPTX skill.

Single-call inspection that returns everything the agent needs:
dimensions, theme fonts, theme colors, color palette, layouts,
slide-by-slide shapes, and table styling.

Usage:
    from pptx_inspect import inspect_template
    inspect_template('template.pptx')
"""

from collections import Counter
from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


# XML namespaces used throughout
_NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


def _emu_to_inches(emu):
    """Convert EMU to inches string."""
    return f"{emu / 914400:.2f}\""


def _extract_theme_fonts(prs):
    """Extract heading and body fonts from the theme XML.

    Accesses theme via slide master relationship (not via element search
    on the master XML, which doesn't contain font definitions).
    """
    try:
        theme_part = prs.slide_masters[0].part.part_related_by(RT.THEME)
        theme_xml = etree.fromstring(theme_part.blob)
        major = theme_xml.findall('.//a:majorFont/a:latin', _NS)
        minor = theme_xml.findall('.//a:minorFont/a:latin', _NS)
        heading = major[0].get('typeface') if major else 'unknown'
        body = minor[0].get('typeface') if minor else 'unknown'
        return heading, body
    except Exception:
        return 'unknown', 'unknown'


def _extract_theme_colors(prs):
    """Extract the theme color scheme (dk1, lt1, dk2, lt2, accent1-6, hlink, folHlink).

    Returns list of (name, hex) tuples.
    """
    try:
        theme_part = prs.slide_masters[0].part.part_related_by(RT.THEME)
        theme_xml = etree.fromstring(theme_part.blob)
        clr_scheme = theme_xml.find('.//a:clrScheme', _NS)
        if clr_scheme is None:
            return []

        colors = []
        # Standard color scheme slots in order
        slots = ['dk1', 'lt1', 'dk2', 'lt2',
                 'accent1', 'accent2', 'accent3', 'accent4',
                 'accent5', 'accent6', 'hlink', 'folHlink']
        for slot in slots:
            el = clr_scheme.find(f'a:{slot}', _NS)
            if el is None:
                continue
            # Color can be srgbClr or sysClr
            srgb = el.find('a:srgbClr', _NS)
            if srgb is not None:
                colors.append((slot, srgb.get('val', '??????')))
                continue
            sys_clr = el.find('a:sysClr', _NS)
            if sys_clr is not None:
                # lastClr is the resolved color
                colors.append((slot, sys_clr.get('lastClr', sys_clr.get('val', '??????'))))
                continue
            colors.append((slot, '??????'))
        return colors
    except Exception:
        return []


def _extract_color_palette(prs):
    """Scan all shapes across all slides for solid fill and text colors.

    Returns Counter of color strings sorted by frequency, capped at 20.
    Includes both #RRGGBB and scheme:name references.
    """
    palette = Counter()

    for slide in prs.slides:
        for shape in slide.shapes:
            _scan_element_colors(shape.element, palette)

    # Sort by frequency, cap at 20
    return palette.most_common(20)


def _scan_element_colors(element, palette):
    """Recursively scan an XML element for color references."""
    # Solid fill sRGB colors
    for srgb in element.findall('.//a:solidFill/a:srgbClr', _NS):
        val = srgb.get('val')
        if val:
            palette[f'#{val}'] += 1

    # Solid fill scheme colors
    for scheme in element.findall('.//a:solidFill/a:schemeClr', _NS):
        val = scheme.get('val')
        if val:
            palette[f'scheme:{val}'] += 1

    # Text run sRGB colors (a:rPr/a:solidFill/a:srgbClr already covered above)
    # Also check for direct color refs on run properties
    for srgb in element.findall('.//a:rPr/a:solidFill/a:srgbClr', _NS):
        val = srgb.get('val')
        if val:
            palette[f'#{val}'] += 1

    for scheme in element.findall('.//a:rPr/a:solidFill/a:schemeClr', _NS):
        val = scheme.get('val')
        if val:
            palette[f'scheme:{val}'] += 1


def _get_fill_info(shape_el):
    """Non-destructive fill detection from XML.

    Does NOT use shape.fill property (which can destructively modify
    inherited backgrounds per _Background.fill docstring).
    Returns a short string describing the fill.
    """
    sp_pr = shape_el.find('.//p:spPr', _NS)
    if sp_pr is None:
        sp_pr = shape_el.find('.//a:spPr', _NS)
    if sp_pr is None:
        # Try alternate paths (group shapes, etc.)
        sp_pr = shape_el

    # Check for solidFill
    solid = sp_pr.find('a:solidFill', _NS)
    if solid is not None:
        srgb = solid.find('a:srgbClr', _NS)
        if srgb is not None:
            return f"fill=#{srgb.get('val', '???')}"
        scheme = solid.find('a:schemeClr', _NS)
        if scheme is not None:
            return f"fill=scheme:{scheme.get('val', '???')}"
        return "fill=SOLID"

    # Check for gradFill
    if sp_pr.find('a:gradFill', _NS) is not None:
        return "fill=GRADIENT"

    # Check for blipFill (image fill)
    if sp_pr.find('a:blipFill', _NS) is not None:
        return "fill=IMAGE"

    # Check for pattFill
    if sp_pr.find('a:pattFill', _NS) is not None:
        return "fill=PATTERN"

    # Check for noFill
    if sp_pr.find('a:noFill', _NS) is not None:
        return "fill=NONE"

    return ""


def _print_layouts(prs):
    """Return string describing all layouts and their placeholders."""
    lines = ["\nLAYOUTS"]
    for i, layout in enumerate(prs.slide_layouts):
        ph_count = len(layout.placeholders)
        lines.append(f"  [{i}] \"{layout.name}\" ({ph_count} placeholder{'s' if ph_count != 1 else ''})")
        for ph in sorted(layout.placeholders, key=lambda p: p.placeholder_format.idx):
            idx = ph.placeholder_format.idx
            ph_type = str(ph.placeholder_format.type).split('(')[0].strip()
            pw = _emu_to_inches(ph.width)
            ph_h = _emu_to_inches(ph.height)
            px = _emu_to_inches(ph.left)
            py = _emu_to_inches(ph.top)
            lines.append(f"    ph[{idx}] {ph_type:<16s} {pw} x {ph_h}  @ ({px}, {py})")
        if ph_count == 0:
            lines.append("    (no placeholders)")
    return "\n".join(lines)


def _print_layout_usage(prs):
    """Return string describing which layouts the template slides actually use."""
    usage = {}  # layout_name -> [slide_numbers]
    for slide_num, slide in enumerate(prs.slides, 1):
        name = slide.slide_layout.name if slide.slide_layout else "unknown"
        usage.setdefault(name, []).append(slide_num)

    all_layout_names = {layout.name for layout in prs.slide_layouts}
    used_names = set(usage.keys())
    unused = all_layout_names - used_names

    lines = ["\nLAYOUT USAGE (which layouts the template slides use)"]
    for name, slides in sorted(usage.items(), key=lambda x: x[1][0]):
        slide_list = ", ".join(str(s) for s in slides)
        lines.append(f"  \"{name}\" → slides: {slide_list}")
    if unused:
        unused_list = ", ".join(f'"{n}"' for n in sorted(unused))
        lines.append(f"  Unused layouts: {unused_list}")
    return "\n".join(lines)


def _print_slide_shapes(prs):
    """Return string with per-slide shape details: name, type, dimensions, position, fill, text preview."""
    lines = ["\nSLIDE-BY-SLIDE SHAPES"]
    for slide_num, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
        lines.append(f"  Slide {slide_num} (layout: \"{layout_name}\"):")
        if len(slide.shapes) == 0:
            lines.append("    (no shapes)")
            continue
        for shape in slide.shapes:
            name = shape.name
            shape_type = str(shape.shape_type).split('(')[0].strip() if shape.shape_type else "UNKNOWN"
            w = _emu_to_inches(shape.width)
            h = _emu_to_inches(shape.height)
            x = _emu_to_inches(shape.left)
            y = _emu_to_inches(shape.top)
            fill = _get_fill_info(shape.element)

            # Text preview (first 50 chars)
            text_preview = ""
            if shape.has_text_frame:
                raw = shape.text_frame.text.strip()
                if raw:
                    preview = raw[:50].replace('\n', ' ')
                    if len(raw) > 50:
                        preview += "..."
                    text_preview = f' text="{preview}"'

            parts = [f'    "{name}" {shape_type} {w}x{h} @ ({x},{y})']
            if fill:
                parts.append(f' {fill}')
            if text_preview:
                parts.append(text_preview)
            lines.append(''.join(parts))
    return "\n".join(lines)


def _print_table_style(prs):
    """Return string with first table's styling details."""
    lines = []
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)
                lines.append(f"\nTABLE STYLE (first table found)")
                lines.append(f"  Location: slide {slide_num}, {rows} rows x {cols} cols")

                # Table element for XML inspection
                tbl = table._tbl

                # Table style ID
                tbl_pr = tbl.find('a:tblPr', _NS)
                if tbl_pr is not None:
                    style_id = tbl_pr.get('bandRow', None)
                    # Active flags
                    flags = []
                    for flag in ('firstRow', 'lastRow', 'firstCol', 'lastCol', 'bandRow', 'bandCol'):
                        val = tbl_pr.get(flag)
                        if val and val == '1':
                            flags.append(flag)
                    if flags:
                        lines.append(f"  Active flags: {', '.join(flags)}")

                    # Style ID from tableStyleId attribute
                    style_el = tbl_pr.find('a:tblStyle', _NS)
                    raw_style_id = tbl_pr.get('tblStyle') or (style_el.text if style_el is not None else None)
                    if not raw_style_id:
                        # Alternate: check for tableStyleId element
                        pass
                    if raw_style_id:
                        lines.append(f"  Table style ID: {raw_style_id}")

                # Sample header row fill
                if rows > 0:
                    lines.append(_print_row_fill(tbl, 0, "Header row"))
                if rows > 1:
                    lines.append(_print_row_fill(tbl, 1, "Data row 1"))

                return "\n".join(lines)  # Only first table
    return ""


def _print_row_fill(tbl, row_idx, label):
    """Return fill info string for cells in a row."""
    row_els = tbl.findall('a:tr', _NS)
    if row_idx >= len(row_els):
        return ""
    row = row_els[row_idx]
    cells = row.findall('a:tc', _NS)
    if not cells:
        return ""

    # Check first cell's fill
    tc = cells[0]
    tc_pr = tc.find('a:tcPr', _NS)
    if tc_pr is None:
        return f"  {label}: (no explicit fill)"

    solid = tc_pr.find('a:solidFill', _NS)
    if solid is not None:
        srgb = solid.find('a:srgbClr', _NS)
        if srgb is not None:
            mods = []
            for child in srgb:
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                val = child.get('val', '')
                mods.append(f"{tag}={val}")
            mod_str = f" ({', '.join(mods)})" if mods else ""
            return f"  {label}: #{srgb.get('val', '???')}{mod_str}"
        scheme = solid.find('a:schemeClr', _NS)
        if scheme is not None:
            mods = []
            for child in scheme:
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                val = child.get('val', '')
                mods.append(f"{tag}={val}")
            mod_str = f" ({', '.join(mods)})" if mods else ""
            return f"  {label}: scheme:{scheme.get('val', '???')}{mod_str}"

    grad = tc_pr.find('a:gradFill', _NS)
    if grad is not None:
        return f"  {label}: GRADIENT"

    return f"  {label}: (no explicit fill)"


def inspect_template(pptx_path):
    """Comprehensive single-call template inspection.

    Returns (and prints) everything the agent needs to build slides:
    dimensions, theme fonts, theme colors, color palette,
    layouts, slide-by-slide shapes, and table styling.
    """
    prs = Presentation(pptx_path)
    lines = []

    # Header
    w = prs.slide_width / 914400
    h = prs.slide_height / 914400
    aspect_label = "16:9" if abs(w / h - 16 / 9) < 0.05 else (
        "4:3" if abs(w / h - 4 / 3) < 0.05 else "custom"
    )
    lines.append(f"Template: {pptx_path}")
    lines.append(f"Dimensions: {w:.2f}\" x {h:.2f}\" ({aspect_label})")
    lines.append(f"Slides: {len(prs.slides)}")

    # Theme fonts
    heading_font, body_font = _extract_theme_fonts(prs)
    lines.append(f"\nTHEME FONTS")
    lines.append(f"  Headings: {heading_font}")
    lines.append(f"  Body: {body_font}")

    # Theme color scheme
    theme_colors = _extract_theme_colors(prs)
    if theme_colors:
        lines.append(f"\nTHEME COLOR SCHEME")
        # Print in rows of 4 for readability
        row = []
        for name, hex_val in theme_colors:
            row.append(f"{name}: {hex_val}")
            if len(row) == 4:
                lines.append(f"  {('  ').join(row)}")
                row = []
        if row:
            lines.append(f"  {('  ').join(row)}")

    # Color palette
    palette = _extract_color_palette(prs)
    if palette:
        lines.append(f"\nCOLOR PALETTE (solid fills + text across all slides, top {len(palette)})")
        for color, count in palette:
            lines.append(f"  {color} — {count} use{'s' if count != 1 else ''}")

    # Layouts
    lines.append(_print_layouts(prs))

    # Layout usage summary
    lines.append(_print_layout_usage(prs))

    # Slide-by-slide shapes
    lines.append(_print_slide_shapes(prs))

    # Table style
    table_style = _print_table_style(prs)
    if table_style:
        lines.append(table_style)

    # Footer
    lines.append(f"\nExisting slides: {len(prs.slides)} (delete after adding new slides)")

    result = "\n".join(lines)
    print(result)
    return result


def _classify_shape(shape):
    """Classify a shape into an editable type.

    Returns (type_str, detail_dict) where type_str is one of:
    TEXT, TEXT_BULLETS, IMAGE, CHART, TABLE, GROUP, OTHER.
    detail_dict contains type-specific current values.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    st = shape.shape_type

    # Table
    if shape.has_table:
        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)
        data = []
        for r in range(rows):
            row_data = []
            for c in range(cols):
                row_data.append(table.cell(r, c).text)
            data.append(row_data)
        return "TABLE", {
            "dimensions": f"{rows} rows x {cols} cols",
            "current_data": data,
            "set_data": None,
            "set_table_from": None,
        }

    # Chart
    if shape.has_chart:
        chart = shape.chart
        chart_type = str(chart.chart_type).split("(")[0].strip() if chart.chart_type else "UNKNOWN"
        detail = {"chart_type": chart_type}
        try:
            plot = chart.plots[0]
            categories = [str(c) for c in plot.categories]
            detail["current_categories"] = categories
            series_dict = {}
            for idx, s in enumerate(plot.series):
                try:
                    name = s.tx.text if hasattr(s, 'tx') and s.tx else f"Series {idx}"
                except Exception:
                    name = f"Series {idx}"
                series_dict[name] = list(s.values)
            detail["current_series"] = series_dict
        except Exception:
            detail["current_categories"] = []
            detail["current_series"] = {}
        detail["set_categories"] = None
        detail["set_series"] = None
        detail["set_chart_from"] = None
        return "CHART", detail

    # Picture / Image
    if st == MSO_SHAPE_TYPE.PICTURE:
        w = f"{shape.width / 914400:.1f}\""
        h = f"{shape.height / 914400:.1f}\""
        return "IMAGE", {
            "current_image": f"image ({w} x {h})",
            "set_image": None,
        }

    # Group
    if st == MSO_SHAPE_TYPE.GROUP:
        child_names = [s.name for s in shape.shapes]
        return "GROUP", {"children": child_names}

    # Text (with bullet detection)
    if shape.has_text_frame:
        tf = shape.text_frame
        paragraphs = [p.text for p in tf.paragraphs]
        non_empty = [p for p in paragraphs if p.strip()]

        # Extract font info from first run
        font_info = {}
        for p in tf.paragraphs:
            if p.runs:
                run = p.runs[0]
                if run.font.name:
                    font_info["current_font_name"] = run.font.name
                if run.font.size:
                    font_info["current_font_size"] = round(run.font.size / 12700, 1)  # EMU to pt
                if run.font.bold is not None:
                    font_info["current_font_bold"] = run.font.bold
                break

        if len(non_empty) > 1:
            detail = {
                "current_text": non_empty,
                "set_text": None,
            }
            detail.update(font_info)
            detail.update({"set_font_name": None, "set_font_size": None, "set_font_bold": None, "set_font_italic": None, "set_font_color": None})
            return "TEXT_BULLETS", detail
        else:
            detail = {
                "current_text": tf.text.strip(),
                "set_text": None,
            }
            detail.update(font_info)
            detail.update({"set_font_name": None, "set_font_size": None, "set_font_bold": None, "set_font_italic": None, "set_font_color": None})
            return "TEXT", detail

    return "OTHER", {}


def _get_slide_colors(slide):
    """Extract hex colors used in a slide for the edit patch template."""
    colors = Counter()
    slide_xml = etree.tostring(slide._element, encoding='unicode')

    import re as _re
    for match in _re.finditer(r'val="([0-9A-Fa-f]{6})"', slide_xml):
        colors[match.group(1).upper()] += 1

    # Return top 10 as dict, skip very common ones like FFFFFF/000000
    result = {}
    for color, count in colors.most_common(15):
        if color not in ('FFFFFF', '000000'):
            result[color] = count
        if len(result) >= 10:
            break
    return result


def _get_available_layouts(prs):
    """Extract available layouts with placeholder names, formatting, and slide usage."""
    # Build usage map: layout_name -> [slide_numbers]
    usage = {}
    for slide_num, slide in enumerate(prs.slides, 1):
        name = slide.slide_layout.name if slide.slide_layout else "unknown"
        usage.setdefault(name, []).append(slide_num)

    layouts = []
    for layout in prs.slide_layouts:
        placeholders = []
        for ph in sorted(layout.placeholders, key=lambda p: p.placeholder_format.idx):
            ph_type = str(ph.placeholder_format.type).split('(')[0].strip()
            ph_info = {
                "name": ph.name,
                "type": ph_type,
            }
            # Extract font formatting from the placeholder's default text
            if ph.has_text_frame and ph.text_frame.paragraphs:
                para = ph.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    f = run.font
                    if f.name:
                        ph_info["font"] = f.name
                    if f.size:
                        ph_info["size"] = f"{f.size.pt:.0f}pt"
                    if f.bold is not None:
                        ph_info["bold"] = f.bold
                    if f.color and f.color.type is not None:
                        try:
                            ph_info["color"] = str(f.color.rgb)
                        except Exception:
                            pass
            placeholders.append(ph_info)

        layout_info = {
            "name": layout.name,
            "placeholders": placeholders,
        }
        used_by = usage.get(layout.name, [])
        if used_by:
            layout_info["used_by_slides"] = used_by
        layouts.append(layout_info)
    return layouts


def generate_edit_patch(pptx_path):
    """Generate a JSON patch template for editing an existing PPTX.

    Returns a dict (also prints it as JSON) with:
    - available_layouts: layouts and their placeholders for adding new slides
    - slides: per-slide shapes with current values and set_* fields (null = no change)
    - structure: slide count, delete/duplicate/reorder controls

    The agent modifies set_* values and passes the result to apply_edit_patch().
    """
    import json

    prs = Presentation(pptx_path)

    patch = {
        "_file": pptx_path,
        "_info": "Modify set_* values below. Delete any slide/shape you don't want to change. Pass to apply_edit_patch().",
        "slide_dimensions": {
            "width_inches": round(prs.slide_width / 914400, 2),
            "height_inches": round(prs.slide_height / 914400, 2),
        },
        "available_layouts": _get_available_layouts(prs),
        "slides": [],
        "add_slides": [],
        "add_slides_from_code": [],
        "structure": {
            "current_slide_count": len(prs.slides),
            "delete_slides": [],
            "reorder": None,
        },
    }

    for slide_num, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
        slide_entry = {
            "slide": slide_num,
            "layout": layout_name,
            "shapes": [],
            "current_colors": _get_slide_colors(slide),
            "color_replace": None,
        }

        for shape in slide.shapes:
            type_str, detail = _classify_shape(shape)
            if type_str == "OTHER":
                continue

            shape_entry = {
                "name": shape.name,
                "type": type_str,
            }
            shape_entry.update(detail)
            slide_entry["shapes"].append(shape_entry)

        if slide_entry["shapes"] or slide_entry["current_colors"]:
            patch["slides"].append(slide_entry)

    output = json.dumps(patch, indent=2, ensure_ascii=False)
    print(output)
    return patch


if __name__ == '__main__':
    import sys
    if len(sys.argv) < 2:
        print("Usage: python pptx_inspect.py <template.pptx>")
        print("       python pptx_inspect.py --edit-patch <template.pptx>")
        sys.exit(1)
    if sys.argv[1] == '--edit-patch' and len(sys.argv) >= 3:
        generate_edit_patch(sys.argv[2])
    else:
        inspect_template(sys.argv[1])
